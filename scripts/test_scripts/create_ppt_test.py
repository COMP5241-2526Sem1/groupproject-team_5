#!/usr/bin/env python3
"""
简单创建PPT测试文件
"""
import sys
sys.path.insert(0, '.venv/Scripts')

from pptx import Presentation

def create_simple_ppt():
    # 创建新演示文稿
    prs = Presentation()
    
    # 添加标题幻灯片
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "人工智能教程"
    subtitle.text = "AI Tutorial - Test Presentation"
    
    # 添加内容幻灯片
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "AI基础概念"
    content.text = "人工智能是计算机科学的重要分支\n主要包括机器学习、深度学习等技术"
    
    # 保存文件
    prs.save('test_ai_presentation.pptx')
    print("PPT文件创建成功")

if __name__ == "__main__":
    create_simple_ppt()