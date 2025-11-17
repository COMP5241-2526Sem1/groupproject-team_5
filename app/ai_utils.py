import os
import re
import requests
import openai
from typing import List, Dict, Any
from collections import Counter
import json
import traceback
import urllib3

# Add file processing support
try:
    from docx import Document
except ImportError:
    Document = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Add dotenv support
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# Add Volcengine SDK support
try:
    from volcenginesdkarkruntime import Ark
except ImportError:
    Ark = None

def create_ark_client(api_key: str, base_url: str = None, timeout: int = 30):
    """
    Create Ark client with SSL verification disabled for Render deployment
    This fixes "getting certificate failed" error on Render
    """
    # Disable SSL warnings for Render deployment
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
    
    # Set environment variable to disable SSL verification (for requests/urllib3)
    os.environ['PYTHONHTTPSVERIFY'] = '0'
    os.environ['CURL_CA_BUNDLE'] = ''
    os.environ['REQUESTS_CA_BUNDLE'] = ''
    
    # Try to create client with httpx (preferred method)
    try:
        import httpx
        http_client = httpx.Client(
            verify=False,
            timeout=timeout
        )
        try:
            # Try with http_client parameter (if SDK supports it)
            if base_url:
                return Ark(api_key=api_key, base_url=base_url, http_client=http_client, timeout=timeout)
            else:
                return Ark(api_key=api_key, http_client=http_client)
        except TypeError:
            # SDK doesn't support http_client parameter, use default client
            # The environment variables should help with SSL issues
            if base_url:
                return Ark(api_key=api_key, base_url=base_url, timeout=timeout)
            else:
                return Ark(api_key=api_key)
    except ImportError:
        # httpx not available, use default client
        # The environment variables should help with SSL issues
        if base_url:
            return Ark(api_key=api_key, base_url=base_url, timeout=timeout)
        else:
            return Ark(api_key=api_key)

def generate_questions(text: str) -> List[str]:
    """Generate questions with enhanced logging"""
    print("=" * 80)
    print(f"🔍 [AI_UTILS] generate_questions() called")
    print(f"   [AI_UTILS] Text length: {len(text)} characters")
    
    # Check for valid API keys in priority order
    ark_api_key = os.environ.get('ARK_API_KEY')
    openai_api_key = os.environ.get('OPENAI_API_KEY')
    
    print(f"   [AI_UTILS] ARK_API_KEY: {'SET (' + str(len(ark_api_key)) + ' chars)' if ark_api_key else 'NOT SET'}")
    print(f"   [AI_UTILS] OPENAI_API_KEY: {'SET (' + str(len(openai_api_key)) + ' chars)' if openai_api_key else 'NOT SET'}")
    
    if ark_api_key and ark_api_key != 'your-bytedance-ark-api-key-here' and len(ark_api_key) > 10:
        print(f"   [AI_UTILS] ✅ Using ARK API (key length: {len(ark_api_key)})")
        print("=" * 80)
        return generate_questions_with_ark(text, ark_api_key)
    elif openai_api_key and openai_api_key != 'your-openai-api-key-here' and openai_api_key.startswith('sk-'):
        print(f"   [AI_UTILS] ✅ Using OpenAI API")
        print("=" * 80)
        return generate_questions_with_openai(text, openai_api_key)
    else:
        print(f"   [AI_UTILS] ⚠️  No valid API key found, using fallback")
        print("=" * 80)
        return generate_questions_fallback(text)

def generate_questions_with_ark(text: str, api_key: str) -> List[str]:
    """Generate questions using ByteDance Ark API with official SDK"""
    if not Ark:
        print("❌ [ARK] Ark SDK not available")
        return generate_questions_fallback(text)
    
    try:
        print(f"🔧 [ARK] Using official ARK SDK method...")
        print(f"   [ARK] API Key: {api_key[:10]}...{api_key[-5:]}")
        print(f"   [ARK] Base URL: https://ark.cn-beijing.volces.com/api/v3")
        
        # Initialize Ark client with SSL verification disabled for Render
        # This fixes "getting certificate failed" error on Render
        client = create_ark_client(api_key)
        
        print(f"📡 [ARK] Calling ARK API with encryption headers...")
        print(f"   [ARK] Model: doubao-1-5-pro-32k-250115")
        print(f"   [ARK] Text length: {len(text)} characters")
        
        # Use official SDK pattern with encryption headers
        completion = client.chat.completions.create(
            model="doubao-1-5-pro-32k-250115",
            messages=[
                {
                    "role": "system",
                    "content": "You are an education expert skilled at generating high-quality classroom interaction questions from teaching text. Please generate 3 questions suitable for classroom interaction based on the given teaching text. Questions should: 1) Test students' understanding of key concepts; 2) Encourage critical thinking; 3) Be suitable for short answer or poll format. Please return 3 questions directly, one per line, without numbering."
                },
                {
                    "role": "user", 
                    "content": f"Please generate 3 classroom interaction questions for the following teaching text:\n\n{text[:2000]}"
                }
            ],
            extra_headers={
                'x-is-encrypted': 'true'  # Official encryption header
            }
        )
        
        print(f"✅ [ARK] ARK API response received with encryption")
        
        questions_text = completion.choices[0].message.content.strip()
        questions = [q.strip() for q in questions_text.split('\n') if q.strip()]
        
        print(f"📝 [ARK] Parsed {len(questions)} questions")
        for i, q in enumerate(questions, 1):
            print(f"   [ARK] Q{i}: {q[:60]}{'...' if len(q) > 60 else ''}")
        
        if len(questions) < 3:
            print(f"⚠️  [ARK] Only {len(questions)} questions generated, adding fallback")
            questions.extend(generate_questions_fallback(text)[:3-len(questions)])
        
        return questions[:3]
        
    except Exception as e:
        print(f"❌ [ARK] Error: {type(e).__name__}: {str(e)[:200]}")
        import traceback
        print(f"   [ARK] Full traceback:")
        traceback.print_exc()
        print(f"   [ARK] Falling back to local questions")
        return generate_questions_fallback(text)

def generate_questions_with_openai(text: str, api_key: str) -> List[str]:
    """Generate questions using OpenAI API"""
    try:
        openai.api_key = api_key
        
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an education expert skilled at generating high-quality classroom interaction questions from teaching text. Please generate 3 questions suitable for classroom interaction based on the given teaching text. Questions should: 1) Test students' understanding of key concepts; 2) Encourage critical thinking; 3) Be suitable for short answer or poll format. Please return 3 questions directly, one per line, without numbering."},
                {"role": "user", "content": f"Please generate 3 classroom interaction questions for the following teaching text:\n\n{text}"}
            ],
            max_tokens=500,
            temperature=0.7
        )
        
        questions_text = response.choices[0].message.content.strip()
        questions = [q.strip() for q in questions_text.split('\n') if q.strip()]
        
        if len(questions) < 3:
            questions.extend(generate_questions_fallback(text)[:3-len(questions)])
        
        return questions[:3]
        
    except Exception as e:
        print(f"OpenAI API error: {e}")
        return generate_questions_fallback(text)

def generate_questions_fallback(text: str) -> List[str]:
    """Improved fallback question generation with better quality"""
    
    # Split sentences
    sentences = re.split(r'[.!?。！？]', text)
    sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
    
    # If text is too short, return generic questions
    if len(sentences) < 2:
        return [
            "What is the main topic discussed in this text?",
            "What key concepts or ideas are presented?",
            "How would you apply this knowledge in practice?"
        ]
    
    questions = []
    
    # Strategy 1: Generate "What is" questions based on first sentence
    first_sentence = sentences[0]
    # Extract topic words (simple method: take first few keywords)
    words = first_sentence.split()[:5]
    subject = ' '.join(words) if len(words) <= 5 else words[0]
    questions.append(f"What is {subject} and why is it important?")
    
    # Strategy 2: Generate questions based on sentence type
    for sentence in sentences[1:min(3, len(sentences))]:
        sentence_lower = sentence.lower()
        
        # Detect definition-type sentences
        if any(word in sentence_lower for word in [' is ', ' are ', ' means ', ' refers to ']):
            # Extract key terms
            key_terms = [w for w in sentence.split() if len(w) > 4][:2]
            if key_terms:
                questions.append(f"Can you explain the relationship between {' and '.join(key_terms)}?")
            else:
                questions.append(f"How would you define the concepts mentioned in: {sentence[:60]}...?")
        
        # Detect function/capability-type sentences  
        elif any(word in sentence_lower for word in ['can ', 'enable', 'allow', 'provide', 'help']):
            questions.append(f"What are the practical applications of: {sentence[:60]}...?")
        
        # Detect process-type sentences
        elif any(word in sentence_lower for word in ['process', 'method', 'approach', 'technique', 'way']):
            questions.append(f"Can you describe how this works: {sentence[:60]}...?")
        
        # Generic questions
        else:
            questions.append(f"What are your thoughts on: {sentence[:60]}...?")
        
        if len(questions) >= 3:
            break
    
    # If still not 3 questions, add critical thinking questions
    if len(questions) < 3:
        critical_questions = [
            "What are the potential limitations or challenges of this approach?",
            "How does this compare to other methods or concepts you know?",
            "What questions do you still have about this topic?"
        ]
        questions.extend(critical_questions[:3 - len(questions)])
    
    return questions[:3]

def generate_activity_from_content(content: str, activity_type: str) -> Dict[str, Any]:
    """Generate a complete activity from teaching content"""
    # Check for valid API keys in priority order
    ark_api_key = os.environ.get('ARK_API_KEY')
    openai_api_key = os.environ.get('OPENAI_API_KEY')
    
    if ark_api_key and ark_api_key != 'your-bytedance-ark-api-key-here' and len(ark_api_key) > 10:
        return generate_activity_with_ark(content, activity_type, ark_api_key)
    elif openai_api_key and openai_api_key != 'your-openai-api-key-here' and openai_api_key.startswith('sk-'):
        return generate_activity_with_openai(content, activity_type, openai_api_key)
    else:
        return generate_activity_fallback(content, activity_type)

def generate_activity_with_ark(content: str, activity_type: str, api_key: str) -> Dict[str, Any]:
    """Generate activity using ByteDance Ark API"""
    if not Ark:
        print("❌ Ark SDK not available")
        return generate_activity_fallback(content, activity_type)
    
    try:
        # Use helper function to create client with SSL verification disabled
        client = create_ark_client(
            api_key=api_key,
            base_url="https://ark.cn-beijing.volces.com/api/v3",
            timeout=30
        )
        
        if activity_type == 'quiz':
            prompt = f"""Based on the following teaching content, create a quiz question with multiple choice options and correct answer.

Content: {content}

Please return a JSON object with:
- title: A descriptive title for the quiz
- question: The quiz question
- options: Array of 4 multiple choice options
- correct_answer: The correct option
- explanation: Brief explanation of why this is correct

Format as valid JSON only."""
        
        elif activity_type == 'poll':
            prompt = f"""Based on the following teaching content, create a poll question with response options.

Content: {content}

Please return a JSON object with:
- title: A descriptive title for the poll
- question: The poll question
- options: Array of 4-6 response options

Format as valid JSON only."""
        
        elif activity_type == 'word_cloud':
            prompt = f"""Based on the following teaching content, create a word cloud activity.

Content: {content}

Please return a JSON object with:
- title: A descriptive title for the word cloud
- question: Instructions for students to submit words/phrases

Format as valid JSON only."""
        
        else:  # short_answer
            prompt = f"""Based on the following teaching content, create a short answer question.

Content: {content}

Please return a JSON object with:
- title: A descriptive title for the question
- question: The short answer question

Format as valid JSON only."""
        
        completion = client.chat.completions.create(
            model="doubao-1-5-pro-32k-250115",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert educator creating interactive learning activities. Always return valid JSON format."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ]
        )
        
        result = json.loads(completion.choices[0].message.content.strip())
        return result
        
    except Exception as e:
        print(f"Ark API error: {e}")
        return generate_activity_fallback(content, activity_type)

def generate_activity_with_openai(content: str, activity_type: str, api_key: str) -> Dict[str, Any]:
    """Generate activity using OpenAI API"""
    try:
        openai.api_key = api_key
        
        if activity_type == 'quiz':
            prompt = f"""Based on the following teaching content, create a quiz question with multiple choice options and correct answer.

Content: {content}

Please return a JSON object with:
- title: A descriptive title for the quiz
- question: The quiz question
- options: Array of 4 multiple choice options
- correct_answer: The correct option
- explanation: Brief explanation of why this is correct

Format as valid JSON only."""
        
        elif activity_type == 'poll':
            prompt = f"""Based on the following teaching content, create a poll question with response options.

Content: {content}

Please return a JSON object with:
- title: A descriptive title for the poll
- question: The poll question
- options: Array of 4-6 response options

Format as valid JSON only."""
        
        elif activity_type == 'word_cloud':
            prompt = f"""Based on the following teaching content, create a word cloud activity.

Content: {content}

Please return a JSON object with:
- title: A descriptive title for the word cloud
- question: Instructions for students to submit words/phrases

Format as valid JSON only."""
        
        else:  # short_answer
            prompt = f"""Based the following teaching content, create a short answer question.

Content: {content}

Please return a JSON object with:
- title: A descriptive title for the question
- question: The short answer question

Format as valid JSON only."""
        
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert educator creating interactive learning activities. Always return valid JSON format."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=500,
            temperature=0.7
        )
        
        result = json.loads(response.choices[0].message.content.strip())
        return result
        
    except Exception as e:
        print(f"OpenAI API error: {e}")
        return generate_activity_fallback(content, activity_type)

def generate_activity_fallback(content: str, activity_type: str) -> Dict[str, Any]:
    """Fallback activity generation without OpenAI"""
    sentences = re.split(r'[.!?。！？]', content)
    sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
    
    if not sentences:
        return {
            'title': 'Generated Activity',
            'question': 'Please discuss the main points of this content.',
            'options': ['Option A', 'Option B', 'Option C', 'Option D'] if activity_type in ['quiz', 'poll'] else None,
            'correct_answer': 'Option A' if activity_type == 'quiz' else None
        }
    
    main_sentence = sentences[0]
    
    if activity_type == 'quiz':
        return {
            'title': 'Understanding Check',
            'question': f'Based on the content, what is the main point about: {main_sentence[:50]}...?',
            'options': ['Option A', 'Option B', 'Option C', 'Option D'],
            'correct_answer': 'Option A',
            'explanation': 'This is the correct answer based on the content.'
        }
    elif activity_type == 'poll':
        return {
            'title': 'Quick Poll',
            'question': f'How well do you understand: {main_sentence[:50]}...?',
            'options': ['Very well', 'Somewhat', 'Not very well', 'Not at all']
        }
    elif activity_type == 'word_cloud':
        return {
            'title': 'Word Association',
            'question': f'What key words come to mind when thinking about: {main_sentence[:50]}...?'
        }
    else:  # short_answer
        return {
            'title': 'Reflection Question',
            'question': f'Please explain your understanding of: {main_sentence[:50]}...'
        }

def group_answers(answers: List[str]) -> Dict[str, Any]:
    """Group and analyze student answers using AI"""
    # Check for valid API keys in priority order
    ark_api_key = os.environ.get('ARK_API_KEY')
    openai_api_key = os.environ.get('OPENAI_API_KEY')
    
    if ark_api_key and ark_api_key != 'your-bytedance-ark-api-key-here' and len(ark_api_key) > 10:
        return group_answers_with_ark(answers, ark_api_key)
    elif openai_api_key and openai_api_key != 'your-openai-api-key-here' and openai_api_key.startswith('sk-'):
        return group_answers_with_openai(answers, openai_api_key)
    else:
        return group_answers_fallback(answers)

def group_answers_with_ark(answers: List[str], api_key: str) -> Dict[str, Any]:
    """Group answers using ByteDance Ark API"""
    if not Ark:
        print("❌ Ark SDK not available")
        return group_answers_fallback(answers)
    
    try:
        # Use helper function to create client with SSL verification disabled
        client = create_ark_client(
            api_key=api_key,
            base_url="https://ark.cn-beijing.volces.com/api/v3",
            timeout=30
        )
        
        answers_text = '\n'.join([f"{i+1}. {answer}" for i, answer in enumerate(answers)])
        
        prompt = f"""Analyze and group the following student answers. Return a JSON object with:
- groups: Array of groups, each with 'name', 'description', 'answers' (array of answer indices)
- summary: Overall summary of the answers
- insights: Key insights from the analysis

Student answers:
{answers_text}

Format as valid JSON only."""
        
        completion = client.chat.completions.create(
            model="doubao-1-5-pro-32k-250115",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert educator analyzing student responses. Group similar answers and provide insights. Always return valid JSON format."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ]
        )
        
        result = json.loads(completion.choices[0].message.content.strip())
        return result
        
    except Exception as e:
        print(f"Ark API error: {e}")
        return group_answers_fallback(answers)

def group_answers_with_openai(answers: List[str], api_key: str) -> Dict[str, Any]:
    """Group answers using OpenAI API"""
    try:
        openai.api_key = api_key
        
        answers_text = '\n'.join([f"{i+1}. {answer}" for i, answer in enumerate(answers)])
        
        prompt = f"""Analyze and group the following student answers. Return a JSON object with:
- groups: Array of groups, each with 'name', 'description', 'answers' (array of answer indices)
- summary: Overall summary of the answers
- insights: Key insights from the analysis

Student answers:
{answers_text}

Format as valid JSON only."""
        
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert educator analyzing student responses. Group similar answers and provide insights. Always return valid JSON format."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=800,
            temperature=0.3
        )
        
        result = json.loads(response.choices[0].message.content.strip())
        return result
        
    except Exception as e:
        print(f"OpenAI API error: {e}")
        return group_answers_fallback(answers)

def group_answers_fallback(answers: List[str]) -> Dict[str, Any]:
    """Fallback answer grouping without OpenAI"""
    # Simple keyword-based grouping
    word_freq = Counter()
    for answer in answers:
        words = re.findall(r'\b\w+\b', answer.lower())
        word_freq.update(words)
    
    # Group by common keywords
    groups = []
    common_words = [word for word, freq in word_freq.most_common(5) if freq > 1]
    
    for i, word in enumerate(common_words[:3]):
        group_answers = [j for j, answer in enumerate(answers) if word in answer.lower()]
        if group_answers:
            groups.append({
                'name': f'Group {i+1}: {word.title()}',
                'description': f'Answers containing "{word}"',
                'answers': group_answers
            })
    
    return {
        'groups': groups,
        'summary': f'Analyzed {len(answers)} responses with {len(groups)} main themes',
        'insights': f'Most common themes: {", ".join(common_words[:3])}'
    }

def extract_text_from_file(file_path: str, file_extension: str) -> str:
    """
    Extract text content from uploaded file
    Supported file formats: .docx, .pdf, .pptx
    """
    try:
        if file_extension.lower() == '.docx':
            return extract_text_from_docx(file_path)
        elif file_extension.lower() == '.pdf':
            return extract_text_from_pdf(file_path)
        elif file_extension.lower() == '.pptx':
            return extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    except Exception as e:
        raise Exception(f"Failed to extract text from file: {str(e)}")

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from Word document"""
    if not Document:
        raise ImportError("python-docx library not installed")
    
    try:
        doc = Document(file_path)
        full_text = []
        
        # Extract paragraph text
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                full_text.append(paragraph.text.strip())
        
        # Extract table text
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        full_text.append(cell.text.strip())
        
        text = '\n'.join(full_text)
        if not text.strip():
            raise ValueError("No text content found in the Word document")
        
        return text
    except Exception as e:
        raise Exception(f"Error reading Word document: {str(e)}")

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file, prefer pdfplumber, fallback to PyPDF2"""
    text = ""
    
    # Prefer pdfplumber as it handles complex PDFs better
    if pdfplumber:
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            if text.strip():
                return text
        except Exception as e:
            print(f"pdfplumber failed: {e}, trying PyPDF2...")
    
    # Fallback to PyPDF2
    if PyPDF2:
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            if text.strip():
                return text
        except Exception as e:
            raise Exception(f"Error reading PDF with PyPDF2: {str(e)}")
    
    if not text.strip():
        raise ValueError("No text content found in the PDF file or PDF libraries not available")
    
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint document"""
    if not Presentation:
        raise ImportError("python-pptx library not installed")
    
    try:
        presentation = Presentation(file_path)
        full_text = []
        
        # Iterate through all slides
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            
            # Extract all text from slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # If it's a table, extract text from table
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_text.append(cell.text.strip())
            
            # If slide has content, add slide title
            if slide_text:
                full_text.append(f"=== Slide {slide_num} ===")
                full_text.extend(slide_text)
                full_text.append("")  # Add blank line separator
        
        text = '\n'.join(full_text)
        if not text.strip():
            raise ValueError("No text content found in the PowerPoint presentation")
        
        return text
    except Exception as e:
        raise Exception(f"Error reading PowerPoint presentation: {str(e)}")

def validate_file_upload(file, allowed_extensions=None):
    """
    Validate uploaded file
    """
    if allowed_extensions is None:
        allowed_extensions = {'.pdf', '.docx', '.pptx'}
    
    if not file or not file.filename:
        return False, "No file selected"
    
    # Check file extension
    file_extension = os.path.splitext(file.filename)[1].lower()
    if file_extension not in allowed_extensions:
        return False, f"Unsupported file format. Allowed formats: {', '.join(allowed_extensions)}"
    
    # Check file size (limit 10MB)
    file.seek(0, 2)  # Move to end of file
    file_size = file.tell()
    file.seek(0)  # Reset file pointer
    
    max_size = 10 * 1024 * 1024  # 10MB
    if file_size > max_size:
        return False, f"File too large. Maximum size allowed: {max_size // (1024*1024)}MB"
    
    return True, "File is valid"
